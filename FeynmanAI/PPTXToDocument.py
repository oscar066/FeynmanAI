import io
from typing import List, Union
from pathlib import Path

from pptx import Presentation
from haystack import Document

class CustomPPTXToDocument:
    def __init__(self):
        try:
            import pptx
        except ImportError:
            raise ImportError("Please install python-pptx: pip install python-pptx")

    def _convert(self, file_content: io.BytesIO) -> str:
        pptx_presentation = Presentation(file_content)
        text_all_slides = []
        for slide in pptx_presentation.slides:
            text_on_slide = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_on_slide.append(shape.text)
            text_all_slides.append("\n".join(text_on_slide))
        return "\f".join(text_all_slides)

    def run(self, sources: List[Union[str, Path]], meta: dict = None):
        documents = []
        for source in sources:
            try:
                with open(source, 'rb') as file:
                    text = self._convert(io.BytesIO(file.read()))
                metadata = meta.copy() if meta else {}
                metadata['source'] = str(source)
                documents.append(Document(content=text, meta=metadata))
            except Exception as e:
                print(f"Error processing {source}: {str(e)}")
        return {"documents": documents}